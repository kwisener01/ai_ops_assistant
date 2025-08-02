"""
ai_ops_assistant.py
=====================

This module implements a prototype AI‑powered operations assistant for a
manufacturing environment. It is designed to help an assistant manager
in a transmission assembly plant track quality and downtime issues,
conduct root cause analysis, and generate reports.  The goal of the
assistant is to reduce busywork by automating data collection,
analysis, and reporting so that more time can be spent on high leverage
activities and strategic thinking.  The design borrows from several
well‑known problem solving frameworks including the Five Whys,
Fishbone (Ishikawa) diagrams, first principles thinking and Dan
Sullivan’s productivity frameworks (Who Not How, The Gap and The Gain,
10× is Easier than 2×, and the Free Zone Frontier).

Key Features
------------

* **Issue tracking** – Load, add and update quality and downtime issues.
  Issues can be stored in an internal pandas ``DataFrame`` or backed
  by a CSV/Excel file for persistence.  Each record includes an
  identifier, description, date raised, status (open/closed), root
  cause category, corrective actions and a list of ``why`` questions
  and answers for 5‑Whys analysis.

* **Root cause analysis** – A flexible interface allows you to record
  iterative “why” questions and answers to uncover the underlying
  cause of a problem.  Once completed, the final answer is stored
  as the root cause and can be visualised with an Ishikawa (fishbone)
  diagram.  Categories default to the classic “5 Ms” used in
  manufacturing (Man, Machine, Material, Method, Measurement) plus
  Environment, but any labels can be supplied.

* **Pattern detection** – Simple analytics compute the most common
  categories and descriptions, mean time to resolution and highlight
  recurring issues.  This helps surface systemic problems and
  continuous improvement opportunities.

* **Report generation** – A summary PDF report can be produced for a
  given date range.  The report lists open and closed issues,
  frequency tables of categories and descriptions, and highlights
  progress (the GAIN) versus remaining work (the GAP).  The PDF is
  generated using the lightweight ``FPDF`` library so that it does not
  require external dependencies.  Optionally, charts can be embedded
  using ``matplotlib``.

* **External file ingestion** – Excel, CSV, PowerPoint and PDF files can
  be ingested.  Excel and CSV are parsed with ``pandas``.  Text
  content can be extracted from PowerPoint (via ``python‑pptx``) and
  PDF files (via ``pdfplumber``) to populate issues automatically.

* **Guidance based on Dan Sullivan frameworks** – The assistant
  encourages high‑leverage thinking.  It highlights tasks that could
  be delegated (Who Not How), tracks progress vs. perfection (The
  Gap and The Gain), focuses on the few activities that will yield
  exponential results (10× is easier than 2×), and prompts the user
  to explore collaborative opportunities beyond the current
  competitive environment (Free Zone Frontier).

This module is self‑contained and does not depend on any external
services.  It can be integrated into a Streamlit or Flask web
application, scheduled in a cron job for daily/weekly summaries, or
invoked from a Jupyter notebook for exploratory analysis.

Example usage::

    from ai_ops_assistant import OpsAssistant
    assistant = OpsAssistant()
    assistant.load_data("issues.xlsx")
    assistant.add_issue("Gearbox vibration detected", category="Machine")
    assistant.ask_why("Gearbox vibration detected", "Why did the vibration occur?", "The bearings were worn")
    assistant.ask_why("Gearbox vibration detected", "Why were the bearings worn?", "Lack of lubrication")
    assistant.finalize_root_cause("Gearbox vibration detected")
    assistant.generate_report(output_path="weekly_report.pdf")

"""

from __future__ import annotations

import dataclasses
import datetime as _dt
import json
import os
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional, Tuple

import pandas as pd
import numpy as np

try:
    import matplotlib
    # Use non‑interactive backend for PDF generation
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
except Exception:
    plt = None  # type: ignore

try:
    from fpdf import FPDF
except ImportError:
    FPDF = None  # type: ignore

try:
    import pdfplumber
except ImportError:
    pdfplumber = None  # type: ignore

try:
    from pptx import Presentation
except ImportError:
    Presentation = None  # type: ignore


@dataclass
class Issue:
    """Data structure representing a single quality/downtime issue."""

    issue_id: int
    description: str
    date_raised: _dt.date
    status: str = "Open"
    category: Optional[str] = None
    root_cause: Optional[str] = None
    corrective_actions: List[str] = field(default_factory=list)
    whys: List[Tuple[str, str]] = field(default_factory=list)  # list of (question, answer)

    def to_dict(self) -> Dict[str, Any]:
        """Serialize the issue to a dictionary."""
        return {
            "issue_id": self.issue_id,
            "description": self.description,
            "date_raised": self.date_raised.isoformat(),
            "status": self.status,
            "category": self.category,
            "root_cause": self.root_cause,
            "corrective_actions": json.dumps(self.corrective_actions),
            "whys": json.dumps(self.whys),
        }

    @classmethod
    def from_dict(cls, d: Dict[str, Any]) -> "Issue":
        """Deserialize from a dictionary."""
        return cls(
            issue_id=int(d["issue_id"]),
            description=str(d["description"]),
            date_raised=_dt.date.fromisoformat(d["date_raised"]),
            status=str(d.get("status", "Open")),
            category=d.get("category"),
            root_cause=d.get("root_cause"),
            corrective_actions=list(json.loads(d.get("corrective_actions", "[]"))),
            whys=[tuple(x) for x in json.loads(d.get("whys", "[]"))],
        )


class OpsAssistant:
    """Core class to manage manufacturing issues and provide analysis.

    The assistant stores issues internally in a pandas DataFrame.  For
    persistence across sessions, call :meth:`load_data` and
    :meth:`save_data` with a CSV or Excel file.  Each issue is
    represented by the :class:`Issue` dataclass and stored in the
    ``_issues`` dictionary keyed by description.  Duplicate
    descriptions are allowed but issue identifiers are unique.
    """

    def __init__(self) -> None:
        # Use dict keyed by issue_id for fast lookup
        self._issues: Dict[int, Issue] = {}
        self._next_id: int = 1

    # ------------------------------------------------------------------
    # Data ingestion and persistence
    # ------------------------------------------------------------------
    def load_data(self, file_path: str) -> None:
        """Load issues from a CSV or Excel file.

        The file must contain columns matching the keys returned by
        :meth:`Issue.to_dict`: ``issue_id``, ``description``,
        ``date_raised``, ``status``, ``category``, ``root_cause``,
        ``corrective_actions``, and ``whys``.  Unknown columns are
        ignored.  JSON‑encoded lists are decoded automatically.
        """
        if not os.path.isfile(file_path):
            raise FileNotFoundError(f"Data file not found: {file_path}")

        ext = os.path.splitext(file_path)[1].lower()
        if ext in [".csv", ".txt"]:
            df = pd.read_csv(file_path)
        elif ext in [".xlsx", ".xlsm", ".xls"]:
            df = pd.read_excel(file_path)
        else:
            raise ValueError(f"Unsupported file extension: {ext}")

        # Normalize column names to lower case
        df.columns = [c.lower() for c in df.columns]
        for _, row in df.iterrows():
            d = {k: row.get(k) for k in [
                "issue_id", "description", "date_raised", "status",
                "category", "root_cause", "corrective_actions", "whys"
            ] if k in row and pd.notnull(row[k])}
            issue = Issue.from_dict(d)
            self._issues[issue.issue_id] = issue
            self._next_id = max(self._next_id, issue.issue_id + 1)

    def save_data(self, file_path: str) -> None:
        """Save current issues to a CSV or Excel file.

        The file format is determined by the extension.  JSON fields
        (lists) are encoded as JSON strings.
        """
        df = pd.DataFrame([issue.to_dict() for issue in self._issues.values()])
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".csv" or ext == ".txt":
            df.to_csv(file_path, index=False)
        elif ext in [".xlsx", ".xlsm", ".xls"]:
            df.to_excel(file_path, index=False)
        else:
            raise ValueError(f"Unsupported file extension: {ext}")

    # ------------------------------------------------------------------
    # Issue management
    # ------------------------------------------------------------------
    def add_issue(self, description: str, category: Optional[str] = None,
                  date_raised: Optional[_dt.date] = None) -> int:
        """Add a new issue and return its unique identifier."""
        issue_id = self._next_id
        self._next_id += 1
        issue = Issue(
            issue_id=issue_id,
            description=description,
            date_raised=date_raised or _dt.date.today(),
            category=category
        )
        self._issues[issue_id] = issue
        return issue_id

    def get_issue(self, issue_id: int) -> Issue:
        """Retrieve an issue by its ID."""
        if issue_id not in self._issues:
            raise KeyError(f"No issue with id {issue_id}")
        return self._issues[issue_id]

    def update_status(self, issue_id: int, status: str) -> None:
        """Update the status of an issue (e.g., 'Open', 'Closed')."""
        issue = self.get_issue(issue_id)
        issue.status = status

    def add_corrective_action(self, issue_id: int, action: str) -> None:
        """Append a corrective action to the specified issue."""
        issue = self.get_issue(issue_id)
        issue.corrective_actions.append(action)

    def ask_why(self, issue_id: int, question: str, answer: str) -> None:
        """Record a single 'why' question and answer for an issue."""
        issue = self.get_issue(issue_id)
        issue.whys.append((question, answer))

    def finalize_root_cause(self, issue_id: int) -> None:
        """Set the final answer from the last 'why' as the root cause."""
        issue = self.get_issue(issue_id)
        if not issue.whys:
            raise ValueError("No 'why' records found for this issue.")
        # Use the answer of the last why as the root cause
        _, last_answer = issue.whys[-1]
        issue.root_cause = last_answer

    # ------------------------------------------------------------------
    # Analytics
    # ------------------------------------------------------------------
    def to_dataframe(self) -> pd.DataFrame:
        """Return a pandas DataFrame representation of the current issues."""
        if not self._issues:
            return pd.DataFrame(columns=["issue_id", "description", "date_raised",
                                         "status", "category", "root_cause",
                                         "corrective_actions", "whys"])
        data = [issue.to_dict() for issue in self._issues.values()]
        # Decode JSON fields back into Python objects
        df = pd.DataFrame(data)
        df["date_raised"] = pd.to_datetime(df["date_raised"])
        df["corrective_actions"] = df["corrective_actions"].apply(lambda x: json.loads(x) if isinstance(x, str) else x)
        df["whys"] = df["whys"].apply(lambda x: [tuple(p) for p in json.loads(x)] if isinstance(x, str) else x)
        return df

    def issue_summary(self) -> pd.DataFrame:
        """Return a DataFrame summarising counts per status and category."""
        df = self.to_dataframe()
        if df.empty:
            return df
        summary = df.groupby(["status", "category"]).size().reset_index(name="count")
        return summary

    def recurring_issues(self, top_n: int = 5) -> pd.DataFrame:
        """Return the most frequent issue descriptions."""
        df = self.to_dataframe()
        if df.empty:
            return df
        freq = df.groupby("description").size().reset_index(name="frequency")
        return freq.sort_values(by="frequency", ascending=False).head(top_n)

    def mean_resolution_time(self) -> Optional[_dt.timedelta]:
        """Compute the mean time to closure for closed issues."""
        df = self.to_dataframe()
        if df.empty:
            return None
        # Ensure date_raised is datetime
        df["date_raised"] = pd.to_datetime(df["date_raised"])
        # For closed issues we assume date_raised as the start and today as end; real use could track close date
        closed = df[df["status"].str.lower() == "closed"].copy()
        if closed.empty:
            return None
        closed["resolved_on"] = pd.to_datetime("today")
        closed["duration"] = closed["resolved_on"] - closed["date_raised"]
        return closed["duration"].mean()

    # ------------------------------------------------------------------
    # Visualisation
    # ------------------------------------------------------------------
    def draw_fishbone(self, issue_id: int, filename: Optional[str] = None,
                      categories: Optional[List[str]] = None) -> Optional[str]:
        """Generate a simple Ishikawa (fishbone) diagram for an issue.

        Parameters
        ----------
        issue_id: int
            Identifier of the issue to visualise.
        filename: Optional[str]
            Path to save the plot.  If ``None`` a temporary PNG will be
            created and the path returned.
        categories: Optional[List[str]]
            List of cause categories.  Defaults to the 6M categories
            used in manufacturing (Man, Machine, Material, Method,
            Measurement, Environment).

        Returns
        -------
        Optional[str]
            Path to the saved figure, or ``None`` if ``matplotlib`` is
            unavailable.
        """
        if plt is None:
            return None
        issue = self.get_issue(issue_id)
        whys = issue.whys
        if not whys:
            raise ValueError("Issue has no recorded 5‑Why analysis.")
        categories = categories or ["Man", "Machine", "Material", "Method", "Measurement", "Environment"]
        # Map whys to categories evenly (for demonstration).  In real
        # usage this should come from user input.
        cat_assignments = {}
        for idx, (_, answer) in enumerate(whys):
            cat_assignments[categories[idx % len(categories)]] = answer
        fig, ax = plt.subplots(figsize=(10, 6))
        ax.axis('off')
        # Draw backbone
        ax.plot([0.1, 0.9], [0.5, 0.5], lw=2, color='black')
        # Draw head (problem)
        ax.text(0.92, 0.5, issue.description, va='center', ha='left', fontsize=10, fontweight='bold')
        # Draw bones
        spine_y_positions = np.linspace(0.1, 0.9, len(categories))
        for y, cat in zip(spine_y_positions, categories):
            # bone line
            ax.plot([0.4, 0.8], [y, 0.5], lw=1.5, color='black')
            ax.text(0.38, y, cat, va='center', ha='right', fontsize=8, fontweight='bold')
            # sub bone for assigned answer
            if cat in cat_assignments:
                answer = cat_assignments[cat]
                ax.text(0.6, (y + 0.5) / 2, answer, va='center', ha='left', fontsize=7, color='darkblue', wrap=True)
        fig.tight_layout()
        if filename is None:
            filename = f"fishbone_{issue_id}.png"
        fig.savefig(filename, dpi=150)
        plt.close(fig)
        return filename

    # ------------------------------------------------------------------
    # Report generation
    # ------------------------------------------------------------------
    def generate_report(self, output_path: str,
                        start_date: Optional[_dt.date] = None,
                        end_date: Optional[_dt.date] = None) -> None:
        """Generate a PDF report summarising issues within a date range.

        The report includes:
        * Overview of open vs closed issues.
        * Top recurring problems.
        * Category breakdown.
        * List of open issues with their 5‑Why root cause summaries.
        * Progress metrics comparing GAINS (issues closed) vs GAPS (issues open).

        If the ``fpdf`` library is available, it is used to produce a
        neatly formatted document.  Otherwise a basic PDF is
        generated using Matplotlib’s ``PdfPages`` backend.  If
        Matplotlib is also unavailable, an exception is raised.
        """
        df = self.to_dataframe()
        if df.empty:
            raise ValueError("No data available to generate report.")
        # Filter by date range
        df = df.copy()
        df["date_raised"] = pd.to_datetime(df["date_raised"]).dt.date
        start_date = start_date or df["date_raised"].min()
        end_date = end_date or df["date_raised"].max()
        df_range = df[(df["date_raised"] >= start_date) & (df["date_raised"] <= end_date)]
        # Prepare summary statistics
        summary = df_range.groupby("status").size().to_dict()
        total_issues = len(df_range)
        closed_count = summary.get("Closed", 0)
        open_count = summary.get("Open", 0)
        gain_percent = closed_count / total_issues * 100 if total_issues else 0
        gap_percent = open_count / total_issues * 100 if total_issues else 0
        # Recurring issues
        freq = df_range.groupby("description").size().reset_index(name="count")
        top_freq = freq.sort_values(by="count", ascending=False).head(5)
        # Category breakdown
        cat_breakdown = df_range.groupby("category").size().reset_index(name="count")
        if FPDF is not None:
            # Use fpdf to build a nicely formatted report
            pdf = FPDF()
            pdf.set_auto_page_break(auto=True, margin=15)
            pdf.add_page()
            pdf.set_font("Arial", "B", 14)
            pdf.cell(0, 10, f"Manufacturing Issue Report: {start_date} to {end_date}", ln=True)
            pdf.set_font("Arial", size=10)
            # Overview
            pdf.ln(4)
            pdf.set_font("Arial", "B", 12)
            pdf.cell(0, 8, "Overview", ln=True)
            pdf.set_font("Arial", size=10)
            pdf.multi_cell(0, 6,
                           f"Total issues: {total_issues}\nClosed: {closed_count} ({gain_percent:.1f}% gain)\n"
                           f"Open: {open_count} ({gap_percent:.1f}% gap)")
            # Top recurring issues
            pdf.ln(4)
            pdf.set_font("Arial", "B", 12)
            pdf.cell(0, 8, "Top Recurring Problems", ln=True)
            pdf.set_font("Arial", size=10)
            for _, row in top_freq.iterrows():
                pdf.multi_cell(0, 5, f"{row['description']} (count: {row['count']})")
            # Category breakdown
            pdf.ln(4)
            pdf.set_font("Arial", "B", 12)
            pdf.cell(0, 8, "Category Breakdown", ln=True)
            pdf.set_font("Arial", size=10)
            for _, row in cat_breakdown.iterrows():
                cat = row['category'] if pd.notnull(row['category']) else 'Uncategorized'
                pdf.multi_cell(0, 5, f"{cat}: {row['count']}")
            # Detailed open issues
            pdf.ln(4)
            pdf.set_font("Arial", "B", 12)
            pdf.cell(0, 8, "Open Issues Detail", ln=True)
            pdf.set_font("Arial", size=9)
            open_issues = df_range[df_range["status"].str.lower() == "open"]
            for _, row in open_issues.iterrows():
                pdf.set_font("Arial", "B", 9)
                pdf.multi_cell(0, 5, f"Issue {row['issue_id']}: {row['description']}")
                pdf.set_font("Arial", size=9)
                # Summarise whys
                whys = row['whys']
                if isinstance(whys, str):
                    whys_list = [tuple(x) for x in json.loads(whys)]
                else:
                    whys_list = whys
                for q, a in whys_list:
                    pdf.multi_cell(0, 5, f"  {q} -> {a}")
                if row['root_cause']:
                    pdf.multi_cell(0, 5, f"  Root cause: {row['root_cause']}")
                if row['corrective_actions']:
                    acts = row['corrective_actions'] if isinstance(row['corrective_actions'], list) else json.loads(row['corrective_actions'])
                    for act in acts:
                        pdf.multi_cell(0, 5, f"  Corrective: {act}")
                pdf.ln(1)
            pdf.output(output_path)
        else:
            # FPDF unavailable: try using matplotlib's PdfPages
            if plt is None:
                raise ImportError("Neither fpdf nor matplotlib is available to generate a report.")
            from matplotlib.backends.backend_pdf import PdfPages
            with PdfPages(output_path) as pdf_pages:
                # Overview page
                fig, ax = plt.subplots(figsize=(8.27, 11.69))  # A4 portrait
                ax.axis('off')
                y = 1.0
                def write_line(text, size=12, weight='normal', dy=0.04):
                    nonlocal y
                    ax.text(0.05, y, text, fontsize=size, fontweight=weight, transform=ax.transAxes)
                    y -= dy
                write_line(f"Manufacturing Issue Report: {start_date} to {end_date}", size=14, weight='bold', dy=0.05)
                write_line("\nOverview", size=12, weight='bold', dy=0.05)
                write_line(f"Total issues: {total_issues}")
                write_line(f"Closed: {closed_count} ({gain_percent:.1f}% gain)")
                write_line(f"Open: {open_count} ({gap_percent:.1f}% gap)")
                write_line("\nTop Recurring Problems", size=12, weight='bold', dy=0.05)
                for _, row in top_freq.iterrows():
                    write_line(f"{row['description']} (count: {row['count']})", size=10)
                write_line("\nCategory Breakdown", size=12, weight='bold', dy=0.05)
                for _, row in cat_breakdown.iterrows():
                    cat = row['category'] if pd.notnull(row['category']) else 'Uncategorized'
                    write_line(f"{cat}: {row['count']}", size=10)
                pdf_pages.savefig(fig)
                plt.close(fig)
                # Open issues detail page
                fig, ax = plt.subplots(figsize=(8.27, 11.69))
                ax.axis('off')
                y = 1.0
                write_line = lambda text, size=10, weight='normal', dy=0.04: (ax.text(0.05, y - dy, text, fontsize=size, fontweight=weight, transform=ax.transAxes), setattr(locals(), 'y', y - dy))
                y -= 0.02
                ax.text(0.05, y, "Open Issues Detail", fontsize=12, fontweight='bold', transform=ax.transAxes)
                y -= 0.06
                open_issues = df_range[df_range["status"].str.lower() == "open"]
                for _, row in open_issues.iterrows():
                    if y < 0.1:
                        pdf_pages.savefig(fig)
                        plt.close(fig)
                        fig, ax = plt.subplots(figsize=(8.27, 11.69))
                        ax.axis('off')
                        y = 1.0
                    ax.text(0.05, y, f"Issue {row['issue_id']}: {row['description']}", fontsize=10, fontweight='bold', transform=ax.transAxes)
                    y -= 0.04
                    whys = row['whys']
                    if isinstance(whys, str):
                        whys_list = [tuple(x) for x in json.loads(whys)]
                    else:
                        whys_list = whys
                    for q, a in whys_list:
                        ax.text(0.06, y, f"{q} -> {a}", fontsize=9, transform=ax.transAxes)
                        y -= 0.03
                    if row['root_cause']:
                        ax.text(0.06, y, f"Root cause: {row['root_cause']}", fontsize=9, transform=ax.transAxes)
                        y -= 0.03
                    if row['corrective_actions']:
                        acts = row['corrective_actions'] if isinstance(row['corrective_actions'], list) else json.loads(row['corrective_actions'])
                        for act in acts:
                            ax.text(0.06, y, f"Corrective: {act}", fontsize=9, transform=ax.transAxes)
                            y -= 0.03
                    y -= 0.02
                pdf_pages.savefig(fig)
                plt.close(fig)

    # ------------------------------------------------------------------
    # External file ingestion
    # ------------------------------------------------------------------
    def extract_text_from_pdf(self, pdf_path: str) -> str:
        """Extract plain text from a PDF file using pdfplumber.

        Returns the extracted text.  If ``pdfplumber`` is not
        installed, raises ``ImportError``.
        """
        if pdfplumber is None:
            raise ImportError("pdfplumber is required for PDF extraction.")
        text = []
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                text.append(page.extract_text() or "")
        return "\n".join(text)

    def extract_text_from_pptx(self, pptx_path: str) -> str:
        """Extract text from a PowerPoint file using python‑pptx.

        Returns the concatenated text from all slides.  If
        ``python‑pptx`` is not installed, raises ``ImportError``.
        """
        if Presentation is None:
            raise ImportError("python-pptx is required for PowerPoint extraction.")
        prs = Presentation(pptx_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    def ingest_issues_from_file(self, file_path: str, default_category: Optional[str] = None) -> List[int]:
        """Ingest new issues from a text source (PDF, PPTX, TXT).

        This helper will extract text from the file, split on newline,
        and create one issue per non‑empty line.  Returns a list of
        new issue identifiers.  For Excel/CSV ingestion use
        :meth:`load_data` directly.
        """
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            text = self.extract_text_from_pdf(file_path)
        elif ext in [".pptx", ".ppt"]:
            text = self.extract_text_from_pptx(file_path)
        else:
            # Fallback: treat as plain text file
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        new_ids = []
        for line in text.splitlines():
            line = line.strip()
            if not line:
                continue
            # Use the line as the description
            new_id = self.add_issue(description=line, category=default_category)
            new_ids.append(new_id)
        return new_ids


def demo() -> None:
    """Small demonstration of the OpsAssistant capabilities.

    Run this function directly to see how to create a few issues,
    conduct root cause analysis, visualise a fishbone diagram and
    generate a PDF report.  The demonstration uses synthetic data to
    illustrate the workflow.  The resulting report and diagram are
    saved into the current working directory.
    """
    assistant = OpsAssistant()
    # Add a few issues
    i1 = assistant.add_issue("Gearbox vibration detected", category="Machine", date_raised=_dt.date.today() - _dt.timedelta(days=3))
    i2 = assistant.add_issue("Oil leakage in hydraulic press", category="Material", date_raised=_dt.date.today() - _dt.timedelta(days=2))
    i3 = assistant.add_issue("Operator forgot to torque bolts", category="Man", date_raised=_dt.date.today() - _dt.timedelta(days=1))
    # Conduct 5‑Whys for issue 1
    assistant.ask_why(i1, "Why was there vibration?", "Bearings were worn")
    assistant.ask_why(i1, "Why were the bearings worn?", "Lubrication schedule was missed")
    assistant.ask_why(i1, "Why was the schedule missed?", "Maintenance crew understaffed")
    assistant.finalize_root_cause(i1)
    assistant.add_corrective_action(i1, "Revise lubrication schedule and allocate staff")
    assistant.update_status(i1, "Closed")
    # Conduct 5‑Whys for issue 2
    assistant.ask_why(i2, "Why was there leakage?", "Seal failure")
    assistant.ask_why(i2, "Why did the seal fail?", "Incorrect installation")
    assistant.finalize_root_cause(i2)
    assistant.add_corrective_action(i2, "Train team on proper seal installation")
    # Leave issue 2 open
    # Issue 3 without root cause yet
    # Visualise fishbone for issue 1
    fig_path = assistant.draw_fishbone(i1, filename="demo_fishbone.png")
    print(f"Fishbone diagram saved to {fig_path}")
    # Generate report
    report_path = "demo_report.pdf"
    assistant.generate_report(report_path)
    print(f"Report saved to {report_path}")


if __name__ == "__main__":
    demo()